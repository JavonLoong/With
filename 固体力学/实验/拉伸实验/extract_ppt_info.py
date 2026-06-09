from pptx import Presentation
import openpyxl
import json

# 1. PPT - 搜索与试件尺寸相关的关键词
ppt = Presentation('拉伸实验PPT-2026年春-1206.pptx')
keywords = ['d0', 'd₀', '直径', '标距', 'L0', 'L₀', '尺寸', '试件', '低碳钢', 'mm', '截面']
results = []
for i, slide in enumerate(ppt.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if any(kw in txt for kw in keywords):
                    results.append(f"Slide {i+1}: {txt}")
        if shape.has_table:
            table = shape.table
            for ri, row in enumerate(table.rows):
                row_data = [cell.text.strip() for cell in row.cells]
                row_text = " | ".join(row_data)
                if any(kw in row_text for kw in keywords):
                    results.append(f"Slide {i+1} Table Row {ri}: {row_text}")

with open('ppt_dimensions.txt', 'w', encoding='utf-8') as f:
    for r in results:
        f.write(r + '\n')

# 2. Excel
wb = openpyxl.load_workbook('拉伸实验数据记录表.xlsx')
with open('excel_dimensions.txt', 'w', encoding='utf-8') as f:
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        f.write(f"--- Sheet: {sheet_name} ---\n")
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=False):
            vals = [str(cell.value) if cell.value is not None else "" for cell in row]
            if any(v for v in vals):
                f.write(" | ".join(vals) + '\n')

print("Done. Files saved.")
