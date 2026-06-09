"""
提取三份参考资料的完整内容，保存为UTF-8文本文件
"""
import sys, io
# 1. DOC文件 (用python-docx但.doc需要另外处理)
# python-docx只能读.docx，对.doc得用其他方式
# 尝试用antiword或直接读取

# 2. PDF文件
import fitz  # PyMuPDF
pdf_path = '拉伸实验指导书2026春.pdf'
doc = fitz.open(pdf_path)
with open('_ref_pdf.txt', 'w', encoding='utf-8') as f:
    for i, page in enumerate(doc):
        text = page.get_text()
        f.write(f"\n{'='*60}\nPage {i+1}\n{'='*60}\n")
        f.write(text)
pc = doc.page_count
doc.close()
print(f"PDF extracted: {pc} pages -> _ref_pdf.txt")

# 3. PPTX文件
from pptx import Presentation
ppt = Presentation('拉伸实验PPT-2026年春-1206.pptx')
with open('_ref_ppt.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(ppt.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        texts.append(txt)
            if shape.has_table:
                table = shape.table
                for ri, row in enumerate(table.rows):
                    row_data = [cell.text.strip() for cell in row.cells]
                    texts.append("TABLE| " + " | ".join(row_data))
        if texts:
            f.write(f"\n--- Slide {i+1} ---\n")
            for t in texts:
                t = t.replace('\uf070', 'π').replace('\uf073', 'σ').replace('\uf065', 'ε')
                f.write(t + '\n')
print(f"PPT extracted: {len(ppt.slides)} slides -> _ref_ppt.txt")

# 4. DOC文件 - 尝试读取
try:
    from docx import Document
    doc2 = Document('拉伸实验数据记录表格.doc')
    with open('_ref_doc.txt', 'w', encoding='utf-8') as f:
        for para in doc2.paragraphs:
            f.write(para.text + '\n')
        for table in doc2.tables:
            f.write('\n--- TABLE ---\n')
            for row in table.rows:
                vals = [cell.text.strip() for cell in row.cells]
                f.write(' | '.join(vals) + '\n')
    print("DOC extracted -> _ref_doc.txt")
except Exception as e:
    print(f"DOC read failed (expected for .doc): {e}")
    # Try alternative: convert to text using COM or just note it
    print("Will try reading as binary text extraction")
    import subprocess
    try:
        result = subprocess.run(['antiword', '拉伸实验数据记录表格.doc'], 
                                capture_output=True, text=True, timeout=10)
        with open('_ref_doc.txt', 'w', encoding='utf-8') as f:
            f.write(result.stdout)
        print("DOC extracted via antiword -> _ref_doc.txt")
    except Exception as e2:
        print(f"antiword also failed: {e2}")
        # Last resort: extract readable text from binary
        with open('拉伸实验数据记录表格.doc', 'rb') as bf:
            raw = bf.read()
        # Try to find text in the binary
        import re
        # Extract strings from DOC binary
        text_parts = []
        try:
            decoded = raw.decode('gbk', errors='ignore')
            # Clean up
            cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', decoded)
            cleaned = re.sub(r'\s+', ' ', cleaned)
            with open('_ref_doc.txt', 'w', encoding='utf-8') as f:
                f.write(cleaned)
            print("DOC extracted via binary decode -> _ref_doc.txt")
        except:
            print("All DOC extraction methods failed")
